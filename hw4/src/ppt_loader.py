import os.path
import uuid
from pathlib import Path
from turtledemo.clock import wochentag
from typing import Iterator

from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTXLoader(BaseLoader):
    def __init__(self, filepath: str, img_dir: str):
        self.filepath = Path(filepath)
        self.img_dir = Path(img_dir)
        if not self.filepath.exists() or not self.filepath.is_file():
            raise FileNotFoundError(f"{filepath} not find.")
        if not self.img_dir.exists():
            self.img_dir.mkdir(exist_ok=True)

    def _extract_text(self, shape):
        text_list = []
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if text.strip():
                    text_list.append(text)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                text = self._extract_text(sub_shape)
                if text.strip():
                    text_list.append(text)
        return "\n".join(filter(None, text_list))

    def _save_img(self, image):
        ext = Path(image.filename).suffix if image.filename else ".png"
        fname = f"{uuid.uuid4().hex}{ext}"
        full_path = os.path.join(self.img_dir, fname)
        with open(full_path, "wb") as f:
            f.write(image.blob)

        return full_path

    def lazy_load(self) -> Iterator[Document]:
        presentation = Presentation(self.filepath.as_posix())
        for page_num, slide in enumerate(presentation.slides, start=1):
            text_list = []
            img_list = []
            for shape in slide.shapes:
                text = self._extract_text(shape)
                if text:
                    text_list.append(text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    fname = self._save_img(shape.image)
                    img_list.append(fname)

            content = f"Slide {page_num}\n\n"
            if text_list:
                content += "\n\n".join(text_list) + "\n\n"

            yield Document(
                page_content=content,
                metadata={
                    "page": page_num,
                    "source": str(self.filepath),
                    "has_images": len(img_list) > 0,
                    "image_count": len(img_list)
                }
            )
